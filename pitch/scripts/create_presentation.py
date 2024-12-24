from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "Echo-Venlig"
    subtitle.text = "Revolutionizing Sustainable Fashion in Denmark"

def create_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "The Problem"
    content.text = """• Denmark generates 85,000 tonnes of textile waste annually
    • Only 7% of textiles are currently recycled
    • Limited access to repair services
    • Disconnected sustainable fashion market"""

def create_solution_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Our Solution"
    content.text = """• Integrated sustainable fashion marketplace
    • Connected repair service network
    • AI-powered sustainability assistant
    • Educational content platform"""

def create_market_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Market Opportunity"
    content.text = """• Danish fashion market: €5.1B
    • Sustainable segment: €765M
    • 15.3% annual growth
    • 78% consumers prefer sustainable options"""

def create_business_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Business Model"
    content.text = """• Marketplace commission (8-15%)
    • Repair service booking fees
    • Premium features
    • Educational content monetization"""

def create_traction_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Current Traction"
    content.text = """• Beta platform launched
    • 20+ sustainable brands onboarded
    • 50+ repair services partnered
    • 1,000+ early access users"""

def create_team_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Our Team"
    content.text = """• Founder & CEO: [Your Name]
    • CTO: [Tech Lead]
    • Sustainability Director: [Expert Name]
    • Advisory Board: Industry Leaders"""

def create_financials_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Financial Projections"
    content.text = """• Year 1: DKK 2.5M
    • Year 2: DKK 6M
    • Year 3: DKK 15M
    • Break-even: Month 18"""

def create_investment_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Investment Ask"
    content.text = """• Seeking: DKK 5M
    • Use of Funds:
      - 40% Development
      - 30% Marketing
      - 20% Team
      - 10% Operations"""

def create_contact_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Contact Us"
    content.text = """• [Your Name]
    • Founder & CEO, Echo-Venlig
    • Email: [your.email]@echo-venlig.dk
    • Phone: +45 XX XX XX XX
    • www.echo-venlig.dk"""

def main():
    prs = Presentation()
    
    # Create all slides
    create_title_slide(prs)
    create_problem_slide(prs)
    create_solution_slide(prs)
    create_market_slide(prs)
    create_business_model_slide(prs)
    create_traction_slide(prs)
    create_team_slide(prs)
    create_financials_slide(prs)
    create_investment_slide(prs)
    create_contact_slide(prs)
    
    # Save the presentation
    prs.save('D:/Maryam/myProject/EchoVenlig/echo_venlig/pitch/Echo-Venlig_Pitch_Deck.pptx')

if __name__ == '__main__':
    main()
